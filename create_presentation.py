from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ─── Color Palette ───────────────────────────────────────────────────────────
C_BG_DARK    = RGBColor(0x07, 0x0B, 0x14)   # near-black navy
C_BG_PANEL   = RGBColor(0x0D, 0x1A, 0x2E)   # deep blue panel
C_ACCENT1    = RGBColor(0x00, 0xC2, 0xFF)   # electric cyan
C_ACCENT2    = RGBColor(0x74, 0x5E, 0xFF)   # indigo-violet
C_ACCENT3    = RGBColor(0x00, 0xFF, 0xC8)   # mint green
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY = RGBColor(0xB0, 0xC4, 0xDE)
C_GOLD       = RGBColor(0xFF, 0xD7, 0x00)
C_RED_SOFT   = RGBColor(0xFF, 0x6B, 0x6B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── Utility helpers ─────────────────────────────────────────────────────────

def add_solid_bg(slide, color: RGBColor):
    """Fill entire slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None,
             line_width_pt=1.5, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,   # MSO_SHAPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=24, bold=False, color=C_WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb


def add_para(tf, text, font_size=14, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p


def add_gradient_rect(slide, l, t, w, h, color1: RGBColor, color2: RGBColor):
    """Add a rectangle with a linear gradient fill (left→right)."""
    shape = add_rect(slide, l, t, w, h)
    sp = shape._element
    spPr = sp.find(nsmap.qn('p:spPr'))

    # Build gradient XML
    gradFill = etree.SubElement(spPr, nsmap.qn('a:gradFill'))
    gradFill.set('flip', 'none')
    gradFill.set('rotWithShape', '1')
    gsLst = etree.SubElement(gradFill, nsmap.qn('a:gsLst'))

    def make_stop(pos, rgb: RGBColor):
        gs = etree.SubElement(gsLst, nsmap.qn('a:gs'))
        gs.set('pos', str(pos * 1000))
        srgb = etree.SubElement(gs, nsmap.qn('a:srgbClr'))
        srgb.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')

    make_stop(0, color1)
    make_stop(100, color2)

    lin = etree.SubElement(gradFill, nsmap.qn('a:lin'))
    lin.set('ang', '5400000')   # 90 degrees (left to right)
    lin.set('scaled', '0')

    # Remove old solidFill if present
    solidFill = spPr.find(nsmap.qn('a:solidFill'))
    if solidFill is not None:
        spPr.remove(solidFill)

    noFill = spPr.find(nsmap.qn('a:noFill'))
    if noFill is not None:
        spPr.remove(noFill)

    return shape


def add_entrance_animation(slide, shape, effect="appear", delay_ms=0, duration_ms=500):
    """
    Add a basic entrance animation to a shape via OOXML timing elements.
    effect: 'appear' | 'fade' | 'fly-in'
    """
    # Get or create timing tree
    spTree = slide.shapes._spTree
    slideEl = spTree.getparent()

    # Find or create <p:timing>
    timing = slideEl.find(nsmap.qn('p:timing'))
    if timing is None:
        timing = etree.SubElement(slideEl, nsmap.qn('p:timing'))

    # Find or create <p:tnLst>
    tnLst = timing.find(nsmap.qn('p:tnLst'))
    if tnLst is None:
        tnLst = etree.SubElement(timing, nsmap.qn('p:tnLst'))

    # Find or create root <p:par>
    rootPar = tnLst.find(nsmap.qn('p:par'))
    if rootPar is None:
        rootPar = etree.SubElement(tnLst, nsmap.qn('p:par'))
        rootCtn = etree.SubElement(rootPar, nsmap.qn('p:cTn'))
        rootCtn.set('id', '1')
        rootCtn.set('dur', 'indefinite')
        rootCtn.set('restart', 'whenNotActive')
        rootCtn.set('nodeType', 'tmRoot')
        childTnLst = etree.SubElement(rootCtn, nsmap.qn('p:childTnLst'))
        bodyPar = etree.SubElement(childTnLst, nsmap.qn('p:par'))
        bodyCtn = etree.SubElement(bodyPar, nsmap.qn('p:cTn'))
        bodyCtn.set('id', '2')
        bodyCtn.set('fill', 'hold')
        bodyCtn.set('nodeType', 'body')
        etree.SubElement(bodyCtn, nsmap.qn('p:childTnLst'))

    # Navigate to body childTnLst
    rootCtn = rootPar.find(nsmap.qn('p:cTn'))
    bodyChildTnLst = rootCtn.find(nsmap.qn('p:childTnLst'))
    bodyPar = bodyChildTnLst.find(nsmap.qn('p:par'))
    if bodyPar is None:
        bodyPar = etree.SubElement(bodyChildTnLst, nsmap.qn('p:par'))
        bodyCtn = etree.SubElement(bodyPar, nsmap.qn('p:cTn'))
        bodyCtn.set('id', '2')
        bodyCtn.set('fill', 'hold')
        bodyCtn.set('nodeType', 'body')
        etree.SubElement(bodyCtn, nsmap.qn('p:childTnLst'))

    bodyCtn = bodyPar.find(nsmap.qn('p:cTn'))
    seq = bodyCtn.find(nsmap.qn('p:childTnLst'))
    if seq is None:
        seq = etree.SubElement(bodyCtn, nsmap.qn('p:childTnLst'))

    # Generate a unique ID for the animation par
    used_ids = [int(el.get('id','0')) for el in timing.iter() if el.get('id')]
    next_id = max(used_ids, default=2) + 1

    shape_id = shape.shape_id

    if effect == "appear":
        preset_id = "1"
        preset_class = "entr"
        accel = "0"
        decel = "0"
        filter_type = "appear"
        sub_type = ""
        anim_tag = "p:set"
    elif effect == "fade":
        preset_id = "10"
        preset_class = "entr"
        filter_type = "fade"
        sub_type = ""
        anim_tag = "p:animEffect"
    else:  # fly-in
        preset_id = "2"
        preset_class = "entr"
        filter_type = "fly"
        sub_type = "from-bottom"
        anim_tag = "p:animEffect"

    par_xml = f"""
<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:cTn id="{next_id}" presetID="{preset_id}" presetClass="{preset_class}"
         presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
    <p:stCondLst>
      <p:cond delay="{delay_ms}"/>
    </p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{next_id+1}" fill="hold">
          <p:stCondLst>
            <p:cond delay="0"/>
          </p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="{next_id+2}" presetID="{preset_id}" presetClass="{preset_class}"
                     presetSubtype="0" fill="hold" grpId="0" nodeType="withEffect">
                <p:stCondLst>
                  <p:cond delay="0"/>
                </p:stCondLst>
                <p:childTnLst>
                  <p:animEffect transition="in" filter="{filter_type}" subtype="{sub_type}">
                    <p:cBhvr>
                      <p:cTn id="{next_id+3}" dur="{duration_ms}" fill="hold"/>
                      <p:tgtEl>
                        <p:spTgt spid="{shape_id}"/>
                      </p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>
"""
    par_el = etree.fromstring(par_xml)
    seq.append(par_el)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank
add_solid_bg(sl, C_BG_DARK)

# Gradient accent bar (top)
bar = add_gradient_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT1, C_ACCENT2)

# Decorative glowing circle (large, blurred feel via semi-transparent rect)
r1 = add_rect(sl, -1.5, 3.5, 6, 6, fill_color=RGBColor(0x00, 0x50, 0x80))
r1.line.fill.background()

# Logo text
logo = add_text_box(sl, "⬡ SOL", 0.4, 0.15, 4, 1.0,
                    font_size=40, bold=True, color=C_ACCENT1, align=PP_ALIGN.LEFT)

# Main headline
h1 = add_text_box(sl, "SOL Data Agent", 0.4, 1.4, 8, 1.6,
                  font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

h2 = add_text_box(sl, "The AI-Powered Autonomous Data Cleaning Platform",
                  0.4, 3.0, 10, 0.9,
                  font_size=24, bold=False, color=C_ACCENT3, align=PP_ALIGN.LEFT)

tagline = add_text_box(sl,
    "Transforming raw, chaotic datasets into production-ready intelligence — in seconds.",
    0.4, 3.85, 10, 0.7,
    font_size=16, bold=False, color=C_LIGHT_GREY, align=PP_ALIGN.LEFT, italic=True)

# Right side — decorative stat chips
chips = [
    ("80%", "of project time\nspent on data prep"),
    ("AI", "Powered Strategist\n4-Phase Pipeline"),
    ("7+", "File Formats\nSupported"),
]
for i, (val, label) in enumerate(chips):
    bx = add_rect(sl, 9.5, 1.3 + i*1.8, 3.4, 1.5,
                  fill_color=C_BG_PANEL, line_color=C_ACCENT2, line_width_pt=1)
    add_text_box(sl, val, 9.6, 1.35 + i*1.8, 1.2, 0.75,
                 font_size=30, bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER)
    add_text_box(sl, label, 10.7, 1.35 + i*1.8, 2.1, 0.75,
                 font_size=11, bold=False, color=C_LIGHT_GREY, align=PP_ALIGN.LEFT)

# Bottom bar
add_gradient_rect(sl, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)

# Slide number / label
add_text_box(sl, "01 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)

# Animations
add_entrance_animation(sl, logo,    effect="fade", delay_ms=0,    duration_ms=600)
add_entrance_animation(sl, h1,      effect="fly-in", delay_ms=300,  duration_ms=700)
add_entrance_animation(sl, h2,      effect="fly-in", delay_ms=700,  duration_ms=600)
add_entrance_animation(sl, tagline, effect="fade", delay_ms=1100, duration_ms=500)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
sl2 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl2, C_BG_DARK)
add_gradient_rect(sl2, 0, 0, 13.33, 0.12, C_RED_SOFT, C_ACCENT2)

sec_label = add_text_box(sl2, "THE PROBLEM", 0.4, 0.2, 4, 0.4,
                         font_size=11, bold=True, color=C_ACCENT1, align=PP_ALIGN.LEFT)

title2 = add_text_box(sl2, "Data Scientists Waste\n80% of Their Time.", 0.4, 0.65, 8, 1.8,
                      font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

problems = [
    ("🗑", "Dirty, incomplete and inconsistent datasets delay every project."),
    ("⏳", "Manual cleaning scripts are fragile, repetitive and error-prone."),
    ("🧩", "Different file formats and SQL engines require bespoke solutions."),
    ("❌", "Human bias in data manipulation leads to silent, costly errors."),
]
for i, (icon, txt) in enumerate(problems):
    bx = add_rect(sl2, 0.4, 2.6 + i*1.05, 12.5, 0.9,
                  fill_color=RGBColor(0x0D, 0x1A, 0x2E), line_color=C_RED_SOFT, line_width_pt=0.8)
    add_text_box(sl2, icon, 0.55, 2.65 + i*1.05, 0.6, 0.7, font_size=22)
    t = add_text_box(sl2, txt, 1.2, 2.68 + i*1.05, 11.4, 0.7,
                     font_size=16, color=C_LIGHT_GREY)
    add_entrance_animation(sl2, bx, effect="fly-in", delay_ms=i*200+300, duration_ms=400)
    add_entrance_animation(sl2, t, effect="fade", delay_ms=i*200+600, duration_ms=300)

add_gradient_rect(sl2, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl2, "02 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)

add_entrance_animation(sl2, sec_label, effect="fade", delay_ms=0, duration_ms=400)
add_entrance_animation(sl2, title2, effect="fly-in", delay_ms=200, duration_ms=600)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION / OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
sl3 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl3, C_BG_DARK)
add_gradient_rect(sl3, 0, 0, 13.33, 0.12, C_ACCENT1, C_ACCENT3)

sec_label3 = add_text_box(sl3, "THE SOLUTION", 0.4, 0.2, 4, 0.4,
                          font_size=11, bold=True, color=C_ACCENT3, align=PP_ALIGN.LEFT)

title3 = add_text_box(sl3, "Meet SOL Data Agent.", 0.4, 0.65, 9, 1.0,
                      font_size=48, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

desc3 = add_text_box(sl3,
    "An autonomous, AI-powered data cleaning platform that transforms messy raw data\n"
    "into production-ready intelligence — without writing a single line of code.",
    0.4, 1.75, 12.5, 1.0,
    font_size=17, color=C_LIGHT_GREY, italic=True)

# 4-phase pipeline diagram
phases = [
    ("01", "ANALYZE", "Scan metadata,\nclassify columns,\ndetect anomalies", C_ACCENT1),
    ("02", "STRATEGIZE", "AI (LLM) crafts\nan optimal cleaning\nplan per dataset", C_ACCENT2),
    ("03", "EXECUTE", "Clean nulls, types,\nformats, outliers &\nduplicates at scale", C_ACCENT3),
    ("04", "VALIDATE", "Side-by-side preview\n+ instant export or\ndirect DB sync", C_GOLD),
]
arrow_x_positions = [0.35, 3.55, 6.75, 9.95]
for i, (num, phase, desc, col) in enumerate(phases):
    x = arrow_x_positions[i]
    box = add_rect(sl3, x, 3.0, 3.0, 3.9,
                   fill_color=C_BG_PANEL, line_color=col, line_width_pt=1.5)
    add_text_box(sl3, num, x+0.1, 3.05, 0.7, 0.55,
                 font_size=28, bold=True, color=col)
    add_text_box(sl3, phase, x+0.1, 3.55, 2.8, 0.65,
                 font_size=15, bold=True, color=C_WHITE)
    add_text_box(sl3, desc, x+0.1, 4.2, 2.8, 1.6,
                 font_size=12, color=C_LIGHT_GREY)
    add_entrance_animation(sl3, box, effect="fly-in", delay_ms=i*250+400, duration_ms=500)

    # Arrow between boxes
    if i < 3:
        arr = add_text_box(sl3, "→", x+3.0, 4.4, 0.3, 0.5,
                           font_size=22, color=C_ACCENT1, align=PP_ALIGN.CENTER)

add_gradient_rect(sl3, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl3, "03 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)

add_entrance_animation(sl3, sec_label3, effect="fade", delay_ms=0, duration_ms=400)
add_entrance_animation(sl3, title3, effect="fly-in", delay_ms=200, duration_ms=600)
add_entrance_animation(sl3, desc3, effect="fade", delay_ms=600, duration_ms=500)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CORE FEATURES (Part 1)
# ═══════════════════════════════════════════════════════════════════════════════
sl4 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl4, C_BG_DARK)
add_gradient_rect(sl4, 0, 0, 13.33, 0.12, C_ACCENT2, C_ACCENT1)

add_text_box(sl4, "CORE FEATURES", 0.4, 0.2, 5, 0.4,
             font_size=11, bold=True, color=C_ACCENT2, align=PP_ALIGN.LEFT)
add_text_box(sl4, "Built for Real-World Data Chaos.", 0.4, 0.65, 10, 0.9,
             font_size=42, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

features_l = [
    ("🗂", "Universal Data Loader",
     "Natively handles CSV, Excel, JSON, XML,\nParquet, Feather, HDF5, ORC — automatically."),
    ("🤖", "AI-Powered Strategist",
     "LLM (Llama-3) reads data context and crafts\na smart Alpha/Beta/Gamma cleaning plan."),
    ("🎯", "User-Defined Goals",
     "Write your cleaning goal in plain English.\nThe AI adapts its entire strategy accordingly."),
]
features_r = [
    ("🔍", "Dynamic Strategy Preview",
     "Inspect every cleaning plan (with JSON schema)\nbefore execution — full transparency."),
    ("📊", "Smart Export Engine",
     "Outputs cleaned data in the exact original\nformat & extension — zero friction."),
    ("🔗", "Direct SQL Connectors",
     "Read, clean and sync data directly\nin SQLite, MySQL, or PostgreSQL databases."),
]

for i, (icon, title, body) in enumerate(features_l):
    bx = add_rect(sl4, 0.35, 1.7 + i*1.8, 6.1, 1.65,
                  fill_color=C_BG_PANEL, line_color=C_ACCENT1, line_width_pt=1)
    add_text_box(sl4, icon + "  " + title, 0.55, 1.75 + i*1.8, 5.7, 0.55,
                 font_size=14, bold=True, color=C_ACCENT1)
    add_text_box(sl4, body, 0.55, 2.25 + i*1.8, 5.7, 0.9,
                 font_size=12, color=C_LIGHT_GREY)
    add_entrance_animation(sl4, bx, effect="fly-in", delay_ms=i*200+300, duration_ms=450)

for i, (icon, title, body) in enumerate(features_r):
    bx = add_rect(sl4, 6.85, 1.7 + i*1.8, 6.1, 1.65,
                  fill_color=C_BG_PANEL, line_color=C_ACCENT2, line_width_pt=1)
    add_text_box(sl4, icon + "  " + title, 7.05, 1.75 + i*1.8, 5.7, 0.55,
                 font_size=14, bold=True, color=C_ACCENT2)
    add_text_box(sl4, body, 7.05, 2.25 + i*1.8, 5.7, 0.9,
                 font_size=12, color=C_LIGHT_GREY)
    add_entrance_animation(sl4, bx, effect="fly-in", delay_ms=i*200+500, duration_ms=450)

add_gradient_rect(sl4, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl4, "04 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INTEGRATED SUBSYSTEMS / TOOLS
# ═══════════════════════════════════════════════════════════════════════════════
sl5 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl5, C_BG_DARK)
add_gradient_rect(sl5, 0, 0, 13.33, 0.12, C_ACCENT3, C_ACCENT2)

add_text_box(sl5, "INTEGRATED SUBSYSTEMS", 0.4, 0.2, 6, 0.4,
             font_size=11, bold=True, color=C_ACCENT3, align=PP_ALIGN.LEFT)
add_text_box(sl5, "A Complete Data Intelligence Ecosystem.", 0.4, 0.65, 12, 0.9,
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

tools = [
    ("⚙", "Semantic Binary Mapping",
     "Auto-detects boolean columns (Yes/No, Approved/Rejected)\nand maps them to 0/1 for ML readiness."),
    ("📷", "OCR Engine",
     "AI-powered text extraction from images and scanned\ndocuments — converts to clean DataFrames."),
    ("📋", "SOL Forms",
     "An advanced data collection system (Google Forms alternative)\nthat prevents dirty data at the source."),
    ("🧪", "AI Lab / Data Corruptor",
     "A sandbox environment to simulate data corruption\nand stress-test cleaning model performance."),
    ("📑", "Audit Report Dashboard",
     "Full transparency reports of every AI intervention,\nchange log, and confidence score per column."),
]

cols = [0.35, 5.0, 9.5]
rows = [1.75, 4.45]
positions = [(cols[0], rows[0]), (cols[1], rows[0]), (cols[2], rows[0]),
             (cols[0], rows[1]), (cols[1], rows[1])]

for i, ((x, y), (icon, title, body)) in enumerate(zip(positions, tools)):
    bx = add_rect(sl5, x, y, 3.8, 2.4,
                  fill_color=C_BG_PANEL, line_color=C_ACCENT3, line_width_pt=1)
    add_text_box(sl5, icon, x+0.1, y+0.1, 0.6, 0.55, font_size=26, color=C_ACCENT3)
    add_text_box(sl5, title, x+0.7, y+0.1, 3.0, 0.55,
                 font_size=13, bold=True, color=C_WHITE)
    add_text_box(sl5, body, x+0.1, y+0.75, 3.55, 1.5,
                 font_size=11, color=C_LIGHT_GREY)
    add_entrance_animation(sl5, bx, effect="fly-in", delay_ms=i*200+300, duration_ms=500)

# Recommendation system note
rec_box = add_rect(sl5, 9.5, rows[1], 3.8, 2.4,
                   fill_color=C_BG_PANEL, line_color=C_GOLD, line_width_pt=1)
add_text_box(sl5, "🎓", 9.6, rows[1]+0.1, 0.6, 0.55, font_size=26, color=C_GOLD)
add_text_box(sl5, "ML Model Advisor", 10.2, rows[1]+0.1, 3.0, 0.55,
             font_size=13, bold=True, color=C_WHITE)
add_text_box(sl5, "Analyzes your cleaned dataset and recommends\nthe optimal ML model architecture\nfor your problem type.", 9.6, rows[1]+0.75, 3.55, 1.5,
             font_size=11, color=C_LIGHT_GREY)
add_entrance_animation(sl5, rec_box, effect="fly-in", delay_ms=1000, duration_ms=500)

add_gradient_rect(sl5, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl5, "05 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════════
sl6 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl6, C_BG_DARK)
add_gradient_rect(sl6, 0, 0, 13.33, 0.12, C_ACCENT1, C_ACCENT2)

add_text_box(sl6, "TECHNOLOGY STACK", 0.4, 0.2, 5, 0.4,
             font_size=11, bold=True, color=C_ACCENT1, align=PP_ALIGN.LEFT)
add_text_box(sl6, "Engineered for Performance & Scalability.", 0.4, 0.65, 11, 0.9,
             font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

tech_groups = [
    ("🧠 AI & Intelligence", C_ACCENT2,
     ["Groq API (Llama-3 8B)", "LangChain / Custom LLM Chains",
      "Dual-Confidence Strategy Engine", "Semantic Column Classifier"]),
    ("⚙ Backend & Core", C_ACCENT1,
     ["FastAPI (REST API Layer)", "Python 3.10+",
      "Pandas / NumPy / SciPy", "SQLAlchemy + PyMySQL + Psycopg2"]),
    ("🎨 Frontend", C_ACCENT3,
     ["Jinja2 HTML Templates", "Custom Sci-Fi Glassmorphism UI",
      "Interactive Preview Components", "Real-time Progress Streaming"]),
    ("📦 Data Layer", C_GOLD,
     ["Universal Loader (7+ formats)", "SQLite / MySQL / PostgreSQL",
      "Parquet / Feather / HDF5 / ORC", "Dynamic Extension Persistence"]),
]

for i, (group, color, items) in enumerate(tech_groups):
    x = 0.35 + (i % 2) * 6.5
    y = 1.75 + (i // 2) * 3.0
    bx = add_rect(sl6, x, y, 6.1, 2.7,
                  fill_color=C_BG_PANEL, line_color=color, line_width_pt=1.5)
    add_text_box(sl6, group, x+0.15, y+0.1, 5.8, 0.5,
                 font_size=15, bold=True, color=color)
    for j, item in enumerate(items):
        add_text_box(sl6, "• " + item, x+0.2, y+0.65+j*0.47, 5.7, 0.4,
                     font_size=12, color=C_LIGHT_GREY)
    add_entrance_animation(sl6, bx, effect="fly-in", delay_ms=i*300+300, duration_ms=500)

add_gradient_rect(sl6, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl6, "06 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BUSINESS MODEL / MONETIZATION
# ═══════════════════════════════════════════════════════════════════════════════
sl7 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl7, C_BG_DARK)
add_gradient_rect(sl7, 0, 0, 13.33, 0.12, C_GOLD, C_ACCENT2)

add_text_box(sl7, "BUSINESS MODEL", 0.4, 0.2, 5, 0.4,
             font_size=11, bold=True, color=C_GOLD, align=PP_ALIGN.LEFT)
add_text_box(sl7, "Four Powerful Revenue Streams.", 0.4, 0.65, 10, 0.9,
             font_size=42, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

models = [
    ("💻", "SaaS Platform", C_ACCENT1,
     "Free: CSVs up to 5MB\nPro: $25/mo — all formats + 1 SQL DB\n"
     "Enterprise: $500/mo — unlimited sync + isolated hosting"),
    ("🏢", "Enterprise Licenses", C_ACCENT2,
     "On-premise installation for enterprises\nwith strict data privacy requirements.\n"
     "Priced per seat · no data leaves their servers"),
    ("⚡", "Supercharged Freelancing", C_ACCENT3,
     "Complete 5-day data cleaning contracts\nin minutes using SOL Agent.\n"
     "Drastic productivity multiplier on Upwork/Fiverr"),
    ("🔌", "API-as-a-Service", C_GOLD,
     "Expose cleaning logic as a REST API.\nCharge per request / token consumed.\n"
     "Plug into any third-party system seamlessly"),
]

for i, (icon, title, color, body) in enumerate(models):
    x = 0.35 + (i % 2) * 6.5
    y = 1.75 + (i // 2) * 2.65
    bx = add_rect(sl7, x, y, 6.1, 2.45,
                  fill_color=C_BG_PANEL, line_color=color, line_width_pt=1.5)
    add_text_box(sl7, icon + "  " + title, x+0.2, y+0.1, 5.7, 0.55,
                 font_size=16, bold=True, color=color)
    add_text_box(sl7, body, x+0.2, y+0.7, 5.7, 1.6,
                 font_size=12, color=C_LIGHT_GREY)
    add_entrance_animation(sl7, bx, effect="fly-in", delay_ms=i*250+300, duration_ms=500)

add_gradient_rect(sl7, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl7, "07 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLOSING / CALL TO ACTION
# ═══════════════════════════════════════════════════════════════════════════════
sl8 = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl8, C_BG_DARK)
add_gradient_rect(sl8, 0, 0, 13.33, 0.12, C_ACCENT1, C_ACCENT3)

# Big glow circle decoration
r_glow = add_rect(sl8, 3.5, 1.0, 6.5, 6.5, fill_color=RGBColor(0x00, 0x20, 0x40))
r_glow.line.fill.background()

add_text_box(sl8, "⬡ SOL Data Agent", 0.5, 0.2, 12.3, 0.7,
             font_size=14, bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER)

close_title = add_text_box(sl8,
    "Stop Cleaning Manually.\nStart Cleaning Intelligently.",
    0.5, 1.3, 12.3, 2.2,
    font_size=46, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

close_sub = add_text_box(sl8,
    "SOL Data Agent — your autonomous AI co-pilot for data quality.",
    0.5, 3.55, 12.3, 0.7,
    font_size=18, italic=True, color=C_ACCENT3, align=PP_ALIGN.CENTER)

# CTA Buttons (visual only)
cta_box = add_rect(sl8, 3.0, 4.55, 3.5, 0.75,
                   fill_color=C_ACCENT1, line_color=C_ACCENT1)
add_text_box(sl8, "🚀  Get Started", 3.0, 4.6, 3.5, 0.65,
             font_size=16, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

cta_box2 = add_rect(sl8, 7.0, 4.55, 3.5, 0.75,
                    fill_color=None, line_color=C_ACCENT1)
add_text_box(sl8, "📩  Request Demo", 7.0, 4.6, 3.5, 0.65,
             font_size=16, bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER)

# Key stats row
stats = [("< 1 min", "Full dataset cleaned"), ("99%", "Bias-free AI decisions"),
         ("7+ formats", "Universal support"), ("4 phases", "Autonomous pipeline")]
for i, (val, label) in enumerate(stats):
    add_text_box(sl8, val, 0.5 + i*3.2, 5.65, 3.0, 0.5,
                 font_size=20, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text_box(sl8, label, 0.5 + i*3.2, 6.1, 3.0, 0.4,
                 font_size=11, color=C_LIGHT_GREY, align=PP_ALIGN.CENTER)

add_gradient_rect(sl8, 0, 7.35, 13.33, 0.15, C_ACCENT2, C_ACCENT1)
add_text_box(sl8, "08 / 08", 12.0, 7.1, 1.2, 0.3,
             font_size=10, color=C_LIGHT_GREY, align=PP_ALIGN.RIGHT)

add_entrance_animation(sl8, close_title, effect="fly-in", delay_ms=300, duration_ms=700)
add_entrance_animation(sl8, close_sub,   effect="fade",   delay_ms=800,  duration_ms=500)
add_entrance_animation(sl8, cta_box,     effect="fly-in", delay_ms=1100, duration_ms=400)
add_entrance_animation(sl8, cta_box2,    effect="fly-in", delay_ms=1300, duration_ms=400)


# ─── Save ────────────────────────────────────────────────────────────────────
OUTPUT = r"e:\run-20260221T125607Z-1-001\run\SOL_Data_Agent_Presentation.pptx"
prs.save(OUTPUT)
print(f"\nSaved -> {OUTPUT}\n")
print("   8 slides  |  Professional animations  |  SOL dark theme")
